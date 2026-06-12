# FILE: app/llm/_pptx_extractor.py
# Purpose: Text extraction from PowerPoint presentations (.pptx).
# Called-by: app.llm.file_analyzer
# Depends-on: stdlib/third-party only
# Last-renovated: 2026-06-11
"""
Text extraction from PowerPoint presentations (.pptx).

Extracts slide text, speaker notes, and table content into a
readable per-slide format suitable for LLM context and RAG indexing.
"""
from __future__ import annotations

import io
import logging
from typing import Optional, Tuple

logger = logging.getLogger(__name__)


def extract_pptx_text(
    file_path: Optional[str] = None,
    file_bytes: Optional[bytes] = None,
) -> Tuple[str, Optional[str]]:
    """
    Extract text content from a PowerPoint file.

    Args:
        file_path: Path to .pptx file on disk.
        file_bytes: Raw file bytes (alternative to path).

    Returns:
        (text, error) — extracted text and optional error message.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "", "python-pptx not installed (pip install python-pptx)"

    try:
        if file_bytes:
            prs = Presentation(io.BytesIO(file_bytes))
        elif file_path:
            prs = Presentation(file_path)
        else:
            return "", "No file provided"
    except Exception as e:
        return "", f"Failed to open presentation: {e}"

    parts: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []

        # Extract text from all shapes
        for shape in slide.shapes:
            # Text frames (titles, body text, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)

            # Tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        slide_parts.append("\t".join(cells))

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"[Speaker Notes] {notes_text}")

        if slide_parts:
            parts.append(f"=== Slide {idx} ===\n" + "\n".join(slide_parts))

    if not parts:
        return "", "Presentation contains no readable text"

    return "\n\n".join(parts), None
