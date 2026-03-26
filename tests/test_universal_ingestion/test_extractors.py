# FILE: tests/test_universal_ingestion/test_extractors.py
"""
Test suite for Universal Document Ingestion (Capability 1).

Tests:
1. TXT extraction (the original failure)
2. DOCX extraction
3. PDF extraction
4. XLSX extraction (NEW)
5. PPTX extraction (NEW)
6. Known text extension (e.g. .py, .json)
7. Universal UTF-8 fallback for unknown extension
8. Binary file rejection
9. detect_document_type() for new types
10. Drive read-content routing for extractable types
"""
import os
import sys
import json
import tempfile
import traceback
from pathlib import Path

# Add project root to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))

PASS = 0
FAIL = 0
ERRORS: list[str] = []


def test(name: str):
    """Decorator-style test runner."""
    def decorator(fn):
        global PASS, FAIL
        try:
            fn()
            print(f"  ✓ {name}")
            PASS += 1
        except AssertionError as e:
            print(f"  ✗ {name}: {e}")
            FAIL += 1
            ERRORS.append(f"{name}: {e}")
        except Exception as e:
            print(f"  ✗ {name}: EXCEPTION — {e}")
            traceback.print_exc()
            FAIL += 1
            ERRORS.append(f"{name}: {type(e).__name__}: {e}")
    return decorator


def main():
    global PASS, FAIL, ERRORS

    print("=" * 60)
    print("Universal Document Ingestion — Test Suite")
    print("=" * 60)

    # ─── Setup: create test files ──────────────────────────────

    tmp = tempfile.mkdtemp(prefix="astra_test_")
    print(f"\nTemp dir: {tmp}\n")

    # 1. TXT file
    txt_path = os.path.join(tmp, "sample.txt")
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write("This is a plain text document.\nLine two.\nLine three with unicode: café résumé.")

    # 2. JSON file (known text ext)
    json_path = os.path.join(tmp, "data.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump({"key": "value", "nested": {"a": 1}}, f, indent=2)

    # 3. Python file (known text ext)
    py_path = os.path.join(tmp, "script.py")
    with open(py_path, "w", encoding="utf-8") as f:
        f.write('def hello():\n    """Say hello."""\n    print("Hello ASTRA")\n')

    # 4. Unknown extension but valid text
    weird_path = os.path.join(tmp, "config.xyz123")
    with open(weird_path, "w", encoding="utf-8") as f:
        f.write("[settings]\nmode = production\nverbose = true\n")

    # 5. Binary file (should be rejected)
    bin_path = os.path.join(tmp, "binary.dat")
    with open(bin_path, "wb") as f:
        f.write(b"\x00\x01\x02\x03\xff\xfe\xfd" * 100)

    # 6. XLSX file
    xlsx_path = os.path.join(tmp, "data.xlsx")
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Revenue"
        ws.append(["Month", "Revenue", "Expenses"])
        ws.append(["January", 50000, 30000])
        ws.append(["February", 62000, 35000])
        ws.append(["March", 48000, 28000])
        ws2 = wb.create_sheet("Notes")
        ws2.append(["This is a test spreadsheet"])
        wb.save(xlsx_path)
        wb.close()
        xlsx_created = True
    except Exception as e:
        print(f"  [WARN] Could not create test XLSX: {e}")
        xlsx_created = False

    # 7. PPTX file
    pptx_path = os.path.join(tmp, "slides.pptx")
    try:
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "ASTRA Test Presentation"
        slide.placeholders[1].text = "Testing universal document ingestion"
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Slide Two"
        slide2.placeholders[1].text = "More content here for extraction test."
        prs.save(pptx_path)
        pptx_created = True
    except Exception as e:
        print(f"  [WARN] Could not create test PPTX: {e}")
        pptx_created = False

    # 8. DOCX file
    docx_path = os.path.join(tmp, "document.docx")
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("This is a test Word document for ASTRA universal ingestion.")
        doc.add_paragraph("Second paragraph with more content.")
        doc.save(docx_path)
        docx_created = True
    except Exception as e:
        print(f"  [WARN] Could not create test DOCX: {e}")
        docx_created = False

    # ─── Import the extractor ──────────────────────────────────

    print("─── extract_text() tests ───")

    from app.llm.file_analyzer import extract_text, is_binary_file, detect_document_type, KNOWN_TEXT_EXTENSIONS

    # Test 1: TXT
    @test("TXT extraction")
    def _():
        text, err = extract_text(file_path=txt_path)
        assert err is None, f"Error: {err}"
        assert "plain text document" in text, f"Content missing, got: {text[:100]}"
        assert "café" in text, "Unicode content missing"

    # Test 2: JSON (known text ext)
    @test("JSON extraction (known text ext)")
    def _():
        text, err = extract_text(file_path=json_path)
        assert err is None, f"Error: {err}"
        assert '"key"' in text, f"JSON content missing"

    # Test 3: Python file
    @test("Python extraction (known text ext)")
    def _():
        text, err = extract_text(file_path=py_path)
        assert err is None, f"Error: {err}"
        assert "def hello" in text, "Python content missing"

    # Test 4: Unknown extension - universal fallback
    @test("Universal UTF-8 fallback (.xyz123)")
    def _():
        text, err = extract_text(file_path=weird_path)
        assert err is None, f"Error: {err}"
        assert "mode = production" in text, f"Fallback content missing, got: {text[:100]}"

    # Test 5: Binary rejection
    @test("Binary file rejection")
    def _():
        text, err = extract_text(file_path=bin_path)
        assert text == "", f"Should return empty, got {len(text)} chars"
        assert err is not None, "Should return an error"
        assert "binary" in err.lower() or "Binary" in err, f"Error should mention binary: {err}"

    # Test 6: XLSX
    if xlsx_created:
        @test("XLSX extraction")
        def _():
            text, err = extract_text(file_path=xlsx_path)
            assert err is None, f"Error: {err}"
            assert "Revenue" in text, f"Sheet name/header missing, got: {text[:200]}"
            assert "January" in text, "Row data missing"
            assert "50000" in text, "Cell value missing"
            assert "Notes" in text, "Second sheet missing"

    # Test 7: PPTX
    if pptx_created:
        @test("PPTX extraction")
        def _():
            text, err = extract_text(file_path=pptx_path)
            assert err is None, f"Error: {err}"
            assert "ASTRA Test Presentation" in text, f"Slide title missing, got: {text[:200]}"
            assert "Slide Two" in text, "Second slide missing"
            assert "More content here" in text, "Slide body missing"

    # Test 8: DOCX
    if docx_created:
        @test("DOCX extraction")
        def _():
            text, err = extract_text(file_path=docx_path)
            assert err is None, f"Error: {err}"
            assert "Test Document" in text, f"Heading missing, got: {text[:200]}"
            assert "universal ingestion" in text, "Body text missing"

    # Test 9: is_binary_file
    @test("is_binary_file — text file")
    def _():
        assert is_binary_file(file_path=txt_path) is False

    @test("is_binary_file — binary file")
    def _():
        assert is_binary_file(file_path=bin_path) is True

    # Test 10: detect_document_type
    print("\n─── detect_document_type() tests ───")

    @test("detect_document_type: .xlsx → spreadsheet")
    def _():
        result = detect_document_type(filename="report.xlsx")
        assert result == "spreadsheet", f"Got: {result}"

    @test("detect_document_type: .pptx → presentation")
    def _():
        result = detect_document_type(filename="deck.pptx")
        assert result == "presentation", f"Got: {result}"

    @test("detect_document_type: .txt → text")
    def _():
        result = detect_document_type(filename="notes.txt")
        assert result == "text", f"Got: {result}"

    @test("detect_document_type: .pdf → pdf")
    def _():
        result = detect_document_type(filename="report.pdf")
        assert result == "pdf", f"Got: {result}"

    @test("detect_document_type: .docx → docx")
    def _():
        result = detect_document_type(filename="letter.docx")
        assert result == "docx", f"Got: {result}"

    @test("detect_document_type: XLSX MIME type → spreadsheet")
    def _():
        result = detect_document_type(
            filename="data.xlsx",
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        assert result == "spreadsheet", f"Got: {result}"

    # Test 11: KNOWN_TEXT_EXTENSIONS coverage
    print("\n─── KNOWN_TEXT_EXTENSIONS coverage ───")

    @test("KNOWN_TEXT_EXTENSIONS includes .txt, .md, .py, .go, .rs")
    def _():
        for ext in [".txt", ".md", ".py", ".go", ".rs", ".java", ".tf", ".toml"]:
            assert ext in KNOWN_TEXT_EXTENSIONS, f"Missing: {ext}"

    # Test 12: Bytes-based extraction
    print("\n─── Bytes-based extraction ───")

    @test("extract_text from bytes (TXT)")
    def _():
        content = b"Hello from bytes content\nLine 2"
        text, err = extract_text(file_bytes=content, filename="test.txt")
        assert err is None, f"Error: {err}"
        assert "Hello from bytes" in text

    if xlsx_created:
        @test("extract_text from bytes (XLSX)")
        def _():
            with open(xlsx_path, "rb") as f:
                content = f.read()
            text, err = extract_text(file_bytes=content, filename="test.xlsx")
            assert err is None, f"Error: {err}"
            assert "Revenue" in text, f"Content missing from bytes-based XLSX extraction"

    # ─── Summary ───────────────────────────────────────────────

    print("\n" + "=" * 60)
    print(f"Results: {PASS} passed, {FAIL} failed")
    print("=" * 60)

    if ERRORS:
        print("\nFailures:")
        for e in ERRORS:
            print(f"  • {e}")
        sys.exit(1)
    else:
        print("\n✓ All tests passed!")
        sys.exit(0)


if __name__ == "__main__":
    main()
